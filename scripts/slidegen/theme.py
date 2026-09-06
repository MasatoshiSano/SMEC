# -*- coding: utf-8 -*-
"""
SMEC スライドテンプレート共通基盤（コンサル系デザイン：紺 + グレー基調）

配色・フォント・スライドサイズ・共通パーツ（タイトルバー、フッター、章番号）を
一元管理する。診断士教科書をスライド化する際は、必ずこのモジュールの関数を通して
図形・テキストを配置し、科目・論点をまたいで見た目の統一感を保つこと。

新しい図解パターンを追加する場合は diagrams.py 側に関数を足し、
色やフォントなどのスタイル値はこのファイルからのみ参照する（直接色コードを
diagrams.py に書かない）。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# スライドサイズ（16:9）
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.5)          # タイトルバー下、本文開始位置
CONTENT_W = Emu(SLIDE_W - MARGIN * 2)  # Length同士の演算はplain intを返すためEmuで包み直す
FOOTER_Y = Inches(7.08)

# ---------------------------------------------------------------------------
# 配色（コンサル系：白背景 + 紺 + グレー、アクセントは抑制的に使う）
# ---------------------------------------------------------------------------
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_NAVY = RGBColor(0x1F, 0x38, 0x64)        # プライマリ（タイトル・見出し・主要枠）
COLOR_NAVY_MID = RGBColor(0x2E, 0x54, 0x95)    # やや明るい紺（図形の塗り）
COLOR_NAVY_LIGHT = RGBColor(0xD9, 0xE2, 0xF3)  # 薄紺（図形の淡い塗り）
COLOR_GRAY_DARK = RGBColor(0x40, 0x40, 0x40)   # 本文テキスト
COLOR_GRAY_MID = RGBColor(0x7F, 0x7F, 0x7F)    # 補助テキスト・罫線
COLOR_GRAY_LIGHT = RGBColor(0xE7, 0xE6, 0xE6)  # 淡いグレー（背景帯・枠）
COLOR_ACCENT = RGBColor(0xC0, 0x00, 0x00)      # 警告・ひっかけポイント用の差し色（暗赤）
COLOR_ACCENT_LIGHT = RGBColor(0xFC, 0xE4, 0xE4)  # 警告ボックスの淡い背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 論点カード用の交互色（過度に使わず、区別が必要な時だけ）
PALETTE_CYCLE = [COLOR_NAVY, RGBColor(0x54, 0x6E, 0x7A), RGBColor(0x8A, 0x9B, 0xA8), RGBColor(0xB0, 0xC4, 0xDE)]

# ---------------------------------------------------------------------------
# フォント（日本語表示を想定。PowerPoint(Windows/Mac)で開く前提のフォント名）
# ---------------------------------------------------------------------------
FONT_TITLE = "Yu Gothic"      # 太字指定と組み合わせて見出しに使う
FONT_BODY = "Yu Gothic"

FS_SLIDE_TITLE = Pt(26)
FS_SECTION_TITLE = Pt(40)
FS_SUBTITLE = Pt(18)
FS_BODY = Pt(16)
FS_SMALL = Pt(12)
FS_LABEL = Pt(14)


def new_presentation():
    """16:9 の空プレゼンテーションを作成する。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    """白背景の白紙スライドを1枚追加して返す。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    return slide


def _set_font(run, size=FS_BODY, color=COLOR_GRAY_DARK, bold=False, font=FONT_BODY, italic=False):
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    # 東アジア言語フォントも明示的に指定（PowerPointの文字化け・フォント差替え防止）
    rPr = run.font._rPr
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_textbox(slide, left, top, width, height, text, size=FS_BODY, color=COLOR_GRAY_DARK,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
                 line_spacing=1.15, italic=False, wrap=True):
    """単一段落のテキストボックスを配置する。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, color=color, bold=bold, font=font, italic=italic)
    return box


def add_bullet_textbox(slide, left, top, width, height, items, size=FS_BODY, color=COLOR_GRAY_DARK,
                        bullet_color=COLOR_NAVY, font=FONT_BODY, line_spacing=1.3, space_after=Pt(8)):
    """箇条書きテキストボックスを配置する。items は str のリスト、または
    (見出しテキスト, インデントレベル0/1) のタプルのリストでも良い。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = space_after
        p.level = level
        marker = "▪ " if level == 0 else "‒ "
        run = p.add_run()
        run.text = marker + text
        _set_font(run, size=size if level == 0 else Pt(size.pt - 2), color=color, font=font,
                   bold=(level == 0 and False))
    return box


def add_title_bar(slide, kicker, title, subject_tag=None):
    """全スライド共通のタイトルバー（上部）を追加する。
    kicker: 小さい上部ラベル（例：'第2部：組織論'）
    title:  スライドのメインタイトル（論点名）
    subject_tag: 右上に出す科目タグ（例：'C. 企業経営理論'）
    """
    # 上部の細い紺色の帯
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    if kicker:
        add_textbox(slide, MARGIN, Inches(0.28), CONTENT_W, Inches(0.3), kicker,
                    size=FS_LABEL, color=COLOR_GRAY_MID, bold=True)

    add_textbox(slide, MARGIN, Inches(0.55), Inches(9.5), Inches(0.7), title,
                size=FS_SLIDE_TITLE, color=COLOR_NAVY, bold=True, font=FONT_TITLE)

    if subject_tag:
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.3), Inches(0.35),
                                      Inches(2.4), Inches(0.42))
        tag.fill.solid()
        tag.fill.fore_color.rgb = COLOR_NAVY_LIGHT
        tag.line.fill.background()
        tag.shadow.inherit = False
        tf = tag.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = 0; tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subject_tag
        _set_font(run, size=Pt(12), color=COLOR_NAVY, bold=True)

    # タイトル下の罫線
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.32), CONTENT_W, Pt(1.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = COLOR_GRAY_LIGHT
    rule.line.fill.background()
    rule.shadow.inherit = False


def add_footer(slide, page_no, total, source_note="出典：中小企業診断士試験 第1次試験 公式シラバス・過去問（SMEC教科書より作成）"):
    """全スライド共通のフッター（下部）を追加する。"""
    add_textbox(slide, MARGIN, FOOTER_Y, Inches(9.5), Inches(0.3), source_note,
                size=Pt(9), color=COLOR_GRAY_MID)
    add_textbox(slide, Inches(12.0), FOOTER_Y, Inches(0.8), Inches(0.3),
                f"{page_no} / {total}", size=Pt(10), color=COLOR_GRAY_MID, align=PP_ALIGN.RIGHT)


def style_box(shape, fill_color=COLOR_NAVY_LIGHT, line_color=None, text_color=COLOR_NAVY,
              radius=None):
    """図形の塗り・枠線・影の共通スタイルを適用する。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_box_with_text(slide, left, top, width, height, text, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                       fill_color=COLOR_NAVY_LIGHT, text_color=COLOR_NAVY, size=FS_BODY, bold=True,
                       align=PP_ALIGN.CENTER, line_color=None, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE,
                       autosize_shrink=True):
    """図形＋中央テキストのセットを1回で配置する（図解パーツの基本単位）。"""
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    style_box(box, fill_color=fill_color, line_color=line_color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, color=text_color, bold=bold, font=font)
    return box


def add_arrow(slide, left, top, width, height, color=COLOR_NAVY, connector=MSO_CONNECTOR.STRAIGHT):
    """2点間を結ぶ矢印コネクタを追加する（プロセスフロー等で使用）。"""
    conn = slide.shapes.add_connector(connector, left, top, left + width, top + height)
    conn.line.color.rgb = color
    conn.line.width = Pt(2.25)
    conn.shadow.inherit = False
    line = conn.line._get_or_add_ln()
    tail = line.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line.append(tail)
    return conn
