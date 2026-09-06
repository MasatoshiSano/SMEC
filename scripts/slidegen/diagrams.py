# -*- coding: utf-8 -*-
"""
SMEC 教科書スライド化のための図解パターン集（12種類）。

各関数は「1枚の完結したスライド」を作る。すべて theme.py の共通パーツ
（タイトルバー・配色・フォント）を通して描画するため、科目・論点が違っても
見た目の統一感が保たれる。

diagrams.py のこの12関数だけを部品として組み合わせれば、
docs/textbook/ の「基本概念 → 具体例 → ひっかけポイント → 過去問で確認する」
という繰り返し構造をスライド化できる、というのがこのテンプレートの狙い。

1.  title_slide          表紙
2.  section_header       章扉
3.  bullet_slide         標準箇条書き（アイコン付き）
4.  concept_definition   用語定義カード（キーワード／定義／具体例）
5.  process_flow         横方向プロセスフロー（手順・ステップ）
6.  matrix_2x2           2×2マトリクス（SWOT・PPM等）
7.  comparison_table     二軸比較表
8.  timeline             時系列・年表
9.  pyramid              階層ピラミッド
10. cycle_diagram        循環図（PDCA等）
11. hub_and_spoke        中心概念＋周辺要素（放射状）
12. caution_box          ひっかけポイント強調ボックス（誤解 vs 正しい理解）
"""

import math
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

from . import theme as th


# ---------------------------------------------------------------------------
# 1. 表紙
# ---------------------------------------------------------------------------
def title_slide(prs, main_title, subtitle, subject_tag):
    slide = th.add_blank_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), th.SLIDE_W, Inches(2.3))
    band.fill.solid()
    band.fill.fore_color.rgb = th.COLOR_NAVY
    band.line.fill.background()
    band.shadow.inherit = False

    th.add_textbox(slide, th.MARGIN, Inches(3.05), th.CONTENT_W, Inches(1.0), main_title,
                    size=Pt(40), color=th.COLOR_WHITE, bold=True, font=th.FONT_TITLE)
    th.add_textbox(slide, th.MARGIN, Inches(3.9), th.CONTENT_W, Inches(0.6), subtitle,
                    size=Pt(18), color=th.COLOR_NAVY_LIGHT)
    th.add_box_with_text(slide, th.MARGIN, Inches(5.3), Inches(3.2), Inches(0.5), subject_tag,
                          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=th.COLOR_NAVY_LIGHT,
                          text_color=th.COLOR_NAVY, size=Pt(14), align=PP_ALIGN.CENTER)
    th.add_textbox(slide, th.MARGIN, Inches(6.9), Inches(6), Inches(0.3),
                    "SMEC - 中小企業診断士試験対策 スライド教材", size=Pt(10), color=th.COLOR_GRAY_MID)
    return slide


# ---------------------------------------------------------------------------
# 2. 章扉
# ---------------------------------------------------------------------------
def section_header(prs, section_no, section_title, description=""):
    slide = th.add_blank_slide(prs)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), th.SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = th.COLOR_NAVY; bar.line.fill.background()
    bar.shadow.inherit = False

    th.add_textbox(slide, Inches(1.0), Inches(2.6), Inches(4), Inches(1.0), section_no,
                    size=Pt(22), color=th.COLOR_GRAY_MID, bold=True)
    th.add_textbox(slide, Inches(1.0), Inches(3.2), Inches(11), Inches(1.2), section_title,
                    size=Pt(36), color=th.COLOR_NAVY, bold=True, font=th.FONT_TITLE)
    if description:
        th.add_textbox(slide, Inches(1.0), Inches(4.3), Inches(10.5), Inches(1.2), description,
                        size=Pt(16), color=th.COLOR_GRAY_DARK, line_spacing=1.3)
    return slide


# ---------------------------------------------------------------------------
# 3. 標準箇条書き
# ---------------------------------------------------------------------------
def bullet_slide(prs, kicker, title, subject_tag, bullets):
    """bullets: [str, ...] または [(str, level), ...]"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)
    th.add_bullet_textbox(slide, th.MARGIN, th.CONTENT_TOP, th.CONTENT_W, Inches(5.2), bullets,
                           size=th.FS_BODY, line_spacing=1.4, space_after=Pt(12))
    return slide


# ---------------------------------------------------------------------------
# 4. 用語定義カード
# ---------------------------------------------------------------------------
def concept_definition(prs, kicker, title, subject_tag, keyword, definition, example):
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    top = th.CONTENT_TOP
    kw_box = th.add_box_with_text(slide, th.MARGIN, top, Inches(3.4), Inches(4.6), keyword,
                                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=th.COLOR_NAVY,
                                   text_color=th.COLOR_WHITE, size=Pt(24), bold=True)
    body_left = th.MARGIN + Inches(3.7)
    body_w = th.CONTENT_W - Inches(3.7)

    th.add_box_with_text(slide, body_left, top, Inches(1.6), Inches(0.4), "定義",
                          shape_type=MSO_SHAPE.RECTANGLE, fill_color=th.COLOR_GRAY_LIGHT,
                          text_color=th.COLOR_NAVY, size=Pt(13), align=PP_ALIGN.LEFT)
    th.add_textbox(slide, body_left, top + Inches(0.55), body_w, Inches(1.8), definition,
                    size=th.FS_BODY, line_spacing=1.35)

    th.add_box_with_text(slide, body_left, top + Inches(2.5), Inches(1.6), Inches(0.4), "具体例",
                          shape_type=MSO_SHAPE.RECTANGLE, fill_color=th.COLOR_GRAY_LIGHT,
                          text_color=th.COLOR_NAVY, size=Pt(13), align=PP_ALIGN.LEFT)
    th.add_textbox(slide, body_left, top + Inches(3.05), body_w, Inches(1.8), example,
                    size=th.FS_BODY, line_spacing=1.35)
    return slide


# ---------------------------------------------------------------------------
# 5. プロセスフロー
# ---------------------------------------------------------------------------
def process_flow(prs, kicker, title, subject_tag, steps):
    """steps: [(label, description), ...]  3〜5ステップ推奨"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    n = len(steps)
    gap = Inches(0.35)
    box_w = (th.CONTENT_W - gap * (n - 1)) / n
    box_h = Inches(1.6)
    top = th.CONTENT_TOP + Inches(1.2)

    for i, (label, desc) in enumerate(steps):
        left = th.MARGIN + i * (box_w + gap)
        num = th.add_box_with_text(slide, left, top - Inches(0.85), Inches(0.5), Inches(0.5), str(i + 1),
                                    shape_type=MSO_SHAPE.OVAL, fill_color=th.COLOR_NAVY,
                                    text_color=th.COLOR_WHITE, size=Pt(16), bold=True)
        th.add_box_with_text(slide, left, top, box_w, box_h, label,
                              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=th.COLOR_NAVY_LIGHT,
                              text_color=th.COLOR_NAVY, size=Pt(15), bold=True)
        if desc:
            th.add_textbox(slide, left, top + box_h + Inches(0.15), box_w, Inches(1.3), desc,
                            size=Pt(12), color=th.COLOR_GRAY_DARK, align=PP_ALIGN.LEFT, line_spacing=1.2)
        if i < n - 1:
            th.add_arrow(slide, left + box_w, top + box_h / 2, gap, 0)
    return slide


# ---------------------------------------------------------------------------
# 6. 2×2マトリクス
# ---------------------------------------------------------------------------
def matrix_2x2(prs, kicker, title, subject_tag, x_label, y_label, quadrants):
    """quadrants: {'tl':(見出し,説明), 'tr':(...), 'bl':(...), 'br':(...)}"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    grid_left = th.MARGIN + Inches(0.9)
    grid_top = th.CONTENT_TOP + Inches(0.2)
    grid_w = th.CONTENT_W - Inches(1.1)
    grid_h = Inches(4.7)
    half_w = grid_w / 2
    half_h = grid_h / 2

    colors = {
        'tl': th.COLOR_NAVY_LIGHT, 'tr': RGBColor(0xC7, 0xD8, 0xEF),
        'bl': th.COLOR_GRAY_LIGHT, 'br': RGBColor(0xE3, 0xEA, 0xF5),
    }
    positions = {
        'tl': (grid_left, grid_top), 'tr': (grid_left + half_w, grid_top),
        'bl': (grid_left, grid_top + half_h), 'br': (grid_left + half_w, grid_top + half_h),
    }
    for key, (label, desc) in quadrants.items():
        left, top = positions[key]
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, half_w, half_h)
        th.style_box(box, fill_color=colors[key], line_color=th.COLOR_WHITE)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Pt(10); tf.margin_top = Pt(10); tf.margin_right = Pt(10)
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = label
        th._set_font(run, size=Pt(16), color=th.COLOR_NAVY, bold=True)
        if desc:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(4)
            run2 = p2.add_run(); run2.text = desc
            th._set_font(run2, size=Pt(12), color=th.COLOR_GRAY_DARK)

    # 軸ラベル
    th.add_textbox(slide, grid_left, grid_top + grid_h + Inches(0.1), grid_w, Inches(0.35),
                    f"→ {x_label}", size=Pt(13), color=th.COLOR_GRAY_MID, align=PP_ALIGN.CENTER, bold=True)
    ylab = slide.shapes.add_textbox(th.MARGIN - Inches(0.15), grid_top, Inches(0.9), grid_h)
    ylab.rotation = -90
    tf = ylab.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f"↑ {y_label}"
    th._set_font(run, size=Pt(13), color=th.COLOR_GRAY_MID, bold=True)
    return slide


# ---------------------------------------------------------------------------
# 7. 二軸比較表
# ---------------------------------------------------------------------------
def comparison_table(prs, kicker, title, subject_tag, col_a_title, col_b_title, rows, row_labels=None):
    """rows: [(a_text, b_text), ...]  row_labels: 各行の見出し（任意）"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    n_rows = len(rows) + 1
    has_label_col = row_labels is not None
    n_cols = 3 if has_label_col else 2

    top = th.CONTENT_TOP + Inches(0.1)
    width = th.CONTENT_W
    height = Inches(0.7) * n_rows if Inches(0.7) * n_rows < Inches(5.0) else Inches(5.0)

    table_shape = slide.shapes.add_table(n_rows, n_cols, th.MARGIN, top, width, height)
    table = table_shape.table

    # 列幅は均等割り＋ラベル列（あれば）だけ狭める
    remain = width - (Inches(1.8) if has_label_col else 0)
    each = remain / 2
    col_idx = 0
    if has_label_col:
        table.columns[0].width = Inches(1.8)
        col_idx = 1
    table.columns[col_idx].width = int(each)
    table.columns[col_idx + 1].width = int(each)

    def style_cell(cell, text, header=False, label=False):
        cell.margin_left = Pt(8); cell.margin_right = Pt(8); cell.margin_top = Pt(6); cell.margin_bottom = Pt(6)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = th.COLOR_NAVY if header else (th.COLOR_GRAY_LIGHT if label else th.COLOR_WHITE)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if header else PP_ALIGN.LEFT
        run = p.add_run(); run.text = text
        th._set_font(run, size=Pt(13) if not header else Pt(14),
                     color=th.COLOR_WHITE if header else th.COLOR_GRAY_DARK, bold=(header or label))

    c0 = 1 if has_label_col else 0
    if has_label_col:
        style_cell(table.cell(0, 0), "", header=True)
    style_cell(table.cell(0, c0), col_a_title, header=True)
    style_cell(table.cell(0, c0 + 1), col_b_title, header=True)

    for r, (a_text, b_text) in enumerate(rows, start=1):
        if has_label_col:
            style_cell(table.cell(r, 0), row_labels[r - 1], label=True)
        style_cell(table.cell(r, c0), a_text)
        style_cell(table.cell(r, c0 + 1), b_text)
    return slide


# ---------------------------------------------------------------------------
# 8. 時系列・年表
# ---------------------------------------------------------------------------
def timeline(prs, kicker, title, subject_tag, milestones):
    """milestones: [(label, desc), ...] 時系列順"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    line_y = th.CONTENT_TOP + Inches(2.0)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, th.MARGIN, line_y, th.CONTENT_W, Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = th.COLOR_NAVY_MID; line.line.fill.background()
    line.shadow.inherit = False

    n = len(milestones)
    step = th.CONTENT_W / n
    for i, (label, desc) in enumerate(milestones):
        cx = th.MARGIN + step * i + step / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), line_y - Inches(0.09), Inches(0.24), Inches(0.24))
        th.style_box(dot, fill_color=th.COLOR_NAVY)
        above = (i % 2 == 0)
        label_top = line_y - Inches(1.5) if above else line_y + Inches(0.4)
        desc_top = line_y - Inches(1.9) if above else line_y + Inches(0.85)
        box_w = step - Inches(0.2)
        th.add_box_with_text(slide, cx - box_w / 2, label_top, box_w, Inches(0.45), label,
                              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=th.COLOR_NAVY_LIGHT,
                              text_color=th.COLOR_NAVY, size=Pt(13), bold=True)
        if desc:
            th.add_textbox(slide, cx - box_w / 2, desc_top, box_w, Inches(1.0), desc,
                            size=Pt(11), color=th.COLOR_GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.15)
    return slide


# ---------------------------------------------------------------------------
# 9. 階層ピラミッド
# ---------------------------------------------------------------------------
def pyramid(prs, kicker, title, subject_tag, levels):
    """levels: 上位から順に [(label, desc), ...]（先頭が頂点）"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    n = len(levels)
    top0 = th.CONTENT_TOP
    total_h = Inches(4.6)
    band_h = total_h / n - Inches(0.08)
    max_w = Inches(6.5)
    min_w = Inches(2.0)
    center_x = th.MARGIN + Inches(3.4)

    for i, (label, desc) in enumerate(levels):
        w = min_w + (max_w - min_w) * (i / max(n - 1, 1))
        top = top0 + i * (band_h + Inches(0.08))
        left = center_x - w / 2
        color = th.COLOR_NAVY if i == 0 else th.PALETTE_CYCLE[min(i, len(th.PALETTE_CYCLE) - 1)]
        shp = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left, top, w, band_h)
        shp.rotation = 180
        th.style_box(shp, fill_color=color, line_color=th.COLOR_WHITE)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = label
        th._set_font(run, size=Pt(14), color=th.COLOR_WHITE, bold=True)

        if desc:
            th.add_textbox(slide, center_x + max_w / 2 + Inches(0.3), top, Inches(3.0), band_h, desc,
                            size=Pt(11), color=th.COLOR_GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    return slide


# ---------------------------------------------------------------------------
# 10. 循環図
# ---------------------------------------------------------------------------
def cycle_diagram(prs, kicker, title, subject_tag, steps):
    """steps: [label, ...] 円状に並べる（3〜6個推奨）"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    n = len(steps)
    cx = th.MARGIN.inches + th.CONTENT_W.inches / 2
    cy = th.CONTENT_TOP.inches + 2.5
    radius = 2.1
    box_w, box_h = 2.0, 0.9

    centers = []
    for i in range(n):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + radius * math.cos(angle)
        y = cy + radius * math.sin(angle)
        centers.append((x, y))

    # 矢印（円弧の代わりに隣接ノード間を直線で結ぶ簡易表現）
    for i in range(n):
        x1, y1 = centers[i]
        x2, y2 = centers[(i + 1) % n]
        # ボックス端に寄せるため少し内側に短縮
        dx, dy = x2 - x1, y2 - y1
        dist = math.hypot(dx, dy)
        shrink = 0.55
        sx = x1 + dx / dist * shrink
        sy = y1 + dy / dist * shrink
        ex = x2 - dx / dist * shrink
        ey = y2 - dy / dist * shrink
        th.add_arrow(slide, Inches(sx), Inches(sy), Inches(ex - sx), Inches(ey - sy))

    for i, label in enumerate(steps):
        x, y = centers[i]
        th.add_box_with_text(slide, Inches(x - box_w / 2), Inches(y - box_h / 2), Inches(box_w), Inches(box_h),
                              label, shape_type=MSO_SHAPE.OVAL, fill_color=th.COLOR_NAVY,
                              text_color=th.COLOR_WHITE, size=Pt(13), bold=True)
    return slide


# ---------------------------------------------------------------------------
# 11. 中心概念＋周辺要素（放射状）
# ---------------------------------------------------------------------------
def hub_and_spoke(prs, kicker, title, subject_tag, hub_text, spokes):
    """spokes: [str, ...]（3〜6個推奨）"""
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    n = len(spokes)
    cx = th.MARGIN.inches + th.CONTENT_W.inches / 2
    cy = th.CONTENT_TOP.inches + 2.6
    radius = 2.6
    hub_r = 1.1
    spoke_w, spoke_h = 2.3, 0.8

    centers = []
    for i in range(n):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + radius * math.cos(angle)
        y = cy + radius * math.sin(angle)
        centers.append((x, y))

    for (x, y) in centers:
        th.add_arrow(slide, Inches(cx), Inches(cy), Inches(x - cx), Inches(y - cy), color=th.COLOR_GRAY_MID)

    th.add_box_with_text(slide, Inches(cx - hub_r), Inches(cy - hub_r), Inches(hub_r * 2), Inches(hub_r * 2),
                          hub_text, shape_type=MSO_SHAPE.OVAL, fill_color=th.COLOR_NAVY,
                          text_color=th.COLOR_WHITE, size=Pt(15), bold=True)

    for (x, y), text in zip(centers, spokes):
        th.add_box_with_text(slide, Inches(x - spoke_w / 2), Inches(y - spoke_h / 2), Inches(spoke_w), Inches(spoke_h),
                              text, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=th.COLOR_NAVY_LIGHT,
                              text_color=th.COLOR_NAVY, size=Pt(12), bold=True)
    return slide


# ---------------------------------------------------------------------------
# 12. ひっかけポイント強調ボックス
# ---------------------------------------------------------------------------
def caution_box(prs, kicker, title, subject_tag, misconception, correct, note=""):
    slide = th.add_blank_slide(prs)
    th.add_title_bar(slide, kicker, title, subject_tag)

    top = th.CONTENT_TOP
    warn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, th.MARGIN, top, th.CONTENT_W, Inches(1.7))
    th.style_box(warn, fill_color=th.COLOR_ACCENT_LIGHT, line_color=th.COLOR_ACCENT)
    tf = warn.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16); tf.margin_top = Pt(10); tf.margin_right = Pt(16)
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = "✗ よくある誤解"
    th._set_font(run, size=Pt(14), color=th.COLOR_ACCENT, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    run2 = p2.add_run(); run2.text = misconception
    th._set_font(run2, size=Pt(15), color=th.COLOR_GRAY_DARK)

    correct_top = top + Inches(2.0)
    ok = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, th.MARGIN, correct_top, th.CONTENT_W, Inches(1.7))
    th.style_box(ok, fill_color=th.COLOR_NAVY_LIGHT, line_color=th.COLOR_NAVY)
    tf2 = ok.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Pt(16); tf2.margin_top = Pt(10); tf2.margin_right = Pt(16)
    p3 = tf2.paragraphs[0]
    run3 = p3.add_run(); run3.text = "✓ 正しい理解"
    th._set_font(run3, size=Pt(14), color=th.COLOR_NAVY, bold=True)
    p4 = tf2.add_paragraph(); p4.space_before = Pt(4)
    run4 = p4.add_run(); run4.text = correct
    th._set_font(run4, size=Pt(15), color=th.COLOR_GRAY_DARK)

    if note:
        th.add_textbox(slide, th.MARGIN, correct_top + Inches(2.0), th.CONTENT_W, Inches(1.0), note,
                        size=Pt(12), color=th.COLOR_GRAY_MID, line_spacing=1.3)
    return slide
