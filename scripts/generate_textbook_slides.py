#!/usr/bin/env python3
"""教科書Markdownから図解中心のPowerPointスライドを生成する。"""

from __future__ import annotations

import argparse
import re
import textwrap
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEXTBOOK_DIR = ROOT / "docs" / "textbook"
OUTPUT_DIR = TEXTBOOK_DIR / "slides"
FONT = "WenQuanYi Micro Hei"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WHITE = "FFFFFF"
INK = "172033"
MUTED = "667085"
PALE = "F5F7FA"
LINE = "D0D5DD"
WARNING = "B42318"
WARNING_BG = "FEF3F2"
EXAMPLE = "175CD3"
EXAMPLE_BG = "EFF8FF"


SUBJECTS = {
    "A": {
        "file": "A_economics_textbook.md",
        "name": "経済学・経済政策",
        "accent": "4F46E5",
        "soft": "EEF2FF",
        "expected": 36,
    },
    "B": {
        "file": "B_finance_accounting_textbook.md",
        "name": "財務・会計",
        "accent": "0F766E",
        "soft": "ECFDF5",
        "expected": 34,
    },
    "C": {
        "file": "C_business_administration_textbook.md",
        "name": "企業経営理論",
        "accent": "C2410C",
        "soft": "FFF7ED",
        "expected": 42,
    },
    "D": {
        "file": "D_operations_management_textbook.md",
        "name": "運営管理",
        "accent": "15803D",
        "soft": "F0FDF4",
        "expected": 35,
    },
    "E": {
        "file": "E_business_law_textbook.md",
        "name": "経営法務",
        "accent": "7E22CE",
        "soft": "FAF5FF",
        "expected": 27,
    },
    "F": {
        "file": "F_information_systems_textbook.md",
        "name": "経営情報システム",
        "accent": "0369A1",
        "soft": "F0F9FF",
        "expected": 24,
    },
    "G": {
        "file": "G_sme_management_policy_textbook.md",
        "name": "中小企業経営・中小企業政策",
        "accent": "BE123C",
        "soft": "FFF1F2",
        "expected": 28,
    },
}


@dataclass
class Topic:
    code: str
    number: int
    title: str
    part: str
    sections: dict[str, list[str]] = field(
        default_factory=lambda: {"core": [], "example": [], "trap": [], "practice": []}
    )


@dataclass
class Textbook:
    subject: str
    title: str
    intro: list[str]
    topics: list[Topic]
    parts: list[str]
    summary: list[str]


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def lighten(hex_color: str, ratio: float = 0.82) -> str:
    values = [int(hex_color[i : i + 2], 16) for i in (0, 2, 4)]
    mixed = [round(v + (255 - v) * ratio) for v in values]
    return "".join(f"{v:02X}" for v in mixed)


def clean_markdown(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^\s*(?:[-*+]|\d+[.)])\s+", "", text)
    text = re.sub(r"^\s*>\s?", "", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"!\[([^\]]*)]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)]\([^)]*\)", r"\1", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"__([^_]+)__", r"\1", text)
    text = re.sub(r"\$+", "", text)
    text = re.sub(r"\\frac\{([^{}]+)}\{([^{}]+)}", r"\1 ÷ \2", text)
    text = re.sub(r"\\text\{([^{}]+)}", r"\1", text)
    replacements = {
        r"\times": "×",
        r"\div": "÷",
        r"\Delta": "Δ",
        r"\alpha": "α",
        r"\beta": "β",
        r"\sigma": "σ",
        r"\sqrt": "√",
        r"\%": "%",
        "&nbsp;": " ",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    text = text.replace("{", "").replace("}", "")
    text = re.sub(r"\s+", " ", text)
    return text.strip(" |")


def shorten(text: str, limit: int) -> str:
    text = clean_markdown(text)
    if len(text) <= limit:
        return text
    return text[: max(1, limit - 1)].rstrip("、。・ ") + "…"


def parse_textbook(path: Path, subject: str) -> Textbook:
    lines = path.read_text(encoding="utf-8").splitlines()
    title = clean_markdown(lines[0].lstrip("# ")) if lines else path.stem
    topic_re = re.compile(rf"^#{{2,3}}\s+({subject})-(\d+)[：:]\s*(.+)$")
    part_re = re.compile(r"^#\s+第\d+部[：:]\s*(.+)$")

    topics: list[Topic] = []
    parts: list[str] = []
    intro: list[str] = []
    summary: list[str] = []
    current_part = "全体"
    current_topic: Topic | None = None
    current_section = "core"
    in_intro = False
    in_summary = False

    for line in lines[1:]:
        stripped = line.strip()
        if stripped.startswith("# まとめ"):
            current_topic = None
            in_summary = True
            in_intro = False
            continue
        if stripped.startswith("## この教科書の使い方"):
            in_intro = True
            in_summary = False
            continue

        part_match = part_re.match(stripped)
        if part_match:
            current_topic = None
            in_intro = False
            in_summary = False
            current_part = clean_markdown(part_match.group(1))
            if current_part not in parts:
                parts.append(current_part)
            continue

        topic_match = topic_re.match(stripped)
        if topic_match:
            in_intro = False
            in_summary = False
            number = int(topic_match.group(2))
            current_topic = Topic(
                code=f"{subject}-{number}",
                number=number,
                title=clean_markdown(topic_match.group(3)),
                part=current_part,
            )
            topics.append(current_topic)
            current_section = "core"
            continue

        if current_topic and stripped.startswith("#"):
            heading = clean_markdown(stripped.lstrip("# "))
            if heading.startswith("基本概念"):
                current_section = "core"
            elif heading.startswith("具体例"):
                current_section = "example"
            elif heading.startswith("ひっかけポイント"):
                current_section = "trap"
            elif heading.startswith("過去問で確認する"):
                current_section = "practice"
            continue

        if current_topic is not None:
            current_topic.sections[current_section].append(line)
        elif in_intro and stripped and stripped != "---":
            intro.append(line)
        elif in_summary and stripped and stripped != "---":
            summary.append(line)

    numbers = [topic.number for topic in topics]
    expected = list(range(1, SUBJECTS[subject]["expected"] + 1))
    if numbers != expected:
        raise ValueError(f"{path.name}: 論点番号が不連続です: {numbers}")
    return Textbook(subject, title, intro, topics, parts, summary)


def markdown_table(lines: Iterable[str]) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in lines:
        stripped = line.strip()
        if not (stripped.startswith("|") and stripped.endswith("|")):
            continue
        cells = [clean_markdown(cell) for cell in stripped.strip("|").split("|")]
        if not any(cells):
            continue
        if all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells):
            continue
        rows.append(cells)
    return rows


def section_text(lines: Iterable[str], limit: int = 140) -> str:
    pieces: list[str] = []
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped == "---" or stripped.startswith("|"):
            continue
        value = clean_markdown(stripped)
        if value and not value.startswith("problem_sets/") and value not in pieces:
            pieces.append(value)
    return shorten(" ".join(pieces), limit)


def split_label_body(text: str) -> tuple[str, str]:
    text = clean_markdown(text)
    for delimiter in ("：", ":"):
        if delimiter in text:
            label, body = text.split(delimiter, 1)
            if 1 <= len(label) <= 24:
                return shorten(label, 22), shorten(body, 60)
    match = re.match(r"(.{2,22}?)(?:とは|は)、?(.+)", text)
    if match:
        return shorten(match.group(1), 22), shorten(match.group(2), 60)
    return shorten(text, 28), ""


def key_items(topic: Topic, maximum: int = 6) -> list[tuple[str, str]]:
    lines = topic.sections["core"]
    candidates: list[tuple[str, str]] = []

    table = markdown_table(lines)
    if len(table) >= 2:
        for row in table[1:]:
            if not row:
                continue
            candidates.append((shorten(row[0], 24), shorten(" / ".join(row[1:]), 62)))

    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith("|") or stripped.startswith("$$"):
            continue
        is_list = bool(re.match(r"^(?:[-*+]|\d+[.)])\s+", stripped))
        bold_match = re.search(r"\*\*([^*]+)\*\*\s*[：:]?\s*(.*)", stripped)
        if bold_match:
            label = clean_markdown(bold_match.group(1))
            body = clean_markdown(bold_match.group(2))
            candidates.append((shorten(label, 24), shorten(body, 62)))
        elif is_list:
            candidates.append(split_label_body(stripped))

    paragraphs: list[str] = []
    buffer: list[str] = []
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith("|") or stripped.startswith("$$"):
            if buffer:
                paragraphs.append(" ".join(buffer))
                buffer = []
            continue
        if not re.match(r"^(?:[-*+]|\d+[.)])\s+", stripped):
            buffer.append(stripped)
    if buffer:
        paragraphs.append(" ".join(buffer))
    for paragraph in paragraphs:
        for sentence in re.split(r"(?<=[。！？])", clean_markdown(paragraph)):
            sentence = sentence.strip()
            if len(sentence) >= 12:
                candidates.append(split_label_body(sentence))

    result: list[tuple[str, str]] = []
    seen: set[str] = set()
    for label, body in candidates:
        key = re.sub(r"\W", "", label)
        if not key or key in seen:
            continue
        seen.add(key)
        result.append((label, body))
        if len(result) >= maximum:
            break

    if not result:
        fallback = section_text(lines, 100)
        result = [(shorten(topic.title, 28), fallback)]
    return result


def formula_text(topic: Topic) -> str:
    formulas: list[str] = []
    in_math = False
    buffer: list[str] = []
    for line in topic.sections["core"] + topic.sections["example"]:
        stripped = line.strip()
        if stripped.startswith("$$"):
            if in_math:
                if buffer:
                    formulas.append(" ".join(buffer))
                buffer = []
                in_math = False
            else:
                in_math = True
            continue
        if in_math:
            buffer.append(stripped)
            continue
        cleaned = clean_markdown(stripped)
        if ("＝" in cleaned or re.search(r"\s=\s", cleaned)) and len(cleaned) <= 100:
            formulas.append(cleaned)
    if buffer:
        formulas.append(" ".join(buffer))
    return shorten(formulas[0], 95) if formulas else ""


def diagram_type(topic: Topic) -> str:
    haystack = topic.title + " " + section_text(topic.sections["core"], 900)
    table = markdown_table(topic.sections["core"])
    formula = formula_text(topic)

    rules = [
        ("matrix", r"PPM|マトリクス|4P|VRIO|PQCDSME|QC7つ|新QC7つ|5S|3R|ゲーム理論|ポートフォリオ"),
        ("pyramid", r"ピラミッド|欲求段階|法階層"),
        ("timeline", r"ライフサイクル|PLC|沿革|変遷|歴史|存続期間|景気動向|短期.*長期|決算処理一巡"),
        ("cycle", r"サイクル|循環|リーン・スタートアップ|アジャイル|JIT|かんばん|PDCA|更新ループ"),
        ("hierarchy", r"階層|OSI|レベル|層|機関設計|ソフトウェア|正規化|分解ツリー"),
        ("network", r"5フォース|クラスター|ネットワーク|SCM|CRM|ステークホルダー|資源依存|情報システムの適用領域"),
        ("comparison", r"対比|比較|vs|VS|種類|形態|類型|分類|方式|地位別|権利|契約|理論|政策"),
        ("process", r"手続|プロセス|フロー|設立|開発|計画|進出|承継|統制|スケジューリング|サプライチェーン"),
    ]
    for visual, pattern in rules:
        if re.search(pattern, haystack, re.IGNORECASE):
            return visual
    if formula:
        return "equation"
    if len(table) >= 3:
        return "table"
    return "cards"


def set_shape_text(
    shape,
    text: str,
    *,
    size: float = 17,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.10,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    **kwargs,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_text(shape, text, **kwargs)
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    fill: str = WHITE,
    line: str = LINE,
    radius: bool = True,
    size: float = 16,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    set_shape_text(
        shape,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
        margin=0.13,
    )
    return shape


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    body: str,
    accent: str,
    soft: str,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.line.color.rgb = rgb(lighten(accent, 0.55))
    shape.line.width = Pt(1.2)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    add_text(
        slide,
        x + 0.22,
        y + 0.10,
        w - 0.35,
        0.34,
        shorten(label, 25),
        size=15,
        color=accent,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    if body:
        add_text(
            slide,
            x + 0.22,
            y + 0.46,
            w - 0.35,
            h - 0.55,
            shorten(body, 72),
            size=12.5,
            color=INK,
            valign=MSO_ANCHOR.TOP,
        )


def add_slide_base(prs: Presentation, accent: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(WHITE)
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.10)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = rgb(accent)
    top_bar.line.fill.background()
    return slide


def add_footer(slide, subject_file: str, number: int, accent: str) -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(7.15), Inches(12.45), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(LINE)
    line.line.fill.background()
    add_text(
        slide,
        0.48,
        7.18,
        10.8,
        0.20,
        f"出典: docs/textbook/{subject_file}",
        size=8.5,
        color=MUTED,
    )
    add_text(
        slide,
        11.75,
        7.16,
        0.9,
        0.22,
        str(number),
        size=9,
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_topic_header(slide, topic: Topic, accent: str) -> None:
    add_box(
        slide,
        0.48,
        0.28,
        0.95,
        0.40,
        topic.code,
        fill=accent,
        line=accent,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        1.58,
        0.25,
        10.95,
        0.48,
        shorten(topic.title, 53),
        size=23,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        10.25,
        0.73,
        2.35,
        0.24,
        shorten(topic.part, 28),
        size=9.5,
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_context_panels(slide, topic: Topic, accent: str) -> None:
    example = section_text(topic.sections["example"], 105)
    trap = section_text(topic.sections["trap"], 105)
    if not trap:
        trap = "用語の定義・適用条件を教科書本文で確認する"

    add_box(
        slide,
        0.55,
        5.83,
        6.03,
        1.07,
        f"具体例｜{example or '身近な事例に置き換えて理解する'}",
        fill=EXAMPLE_BG,
        line=lighten(EXAMPLE, 0.55),
        size=12.5,
        color=EXAMPLE,
        bold=False,
    )
    add_box(
        slide,
        6.75,
        5.83,
        6.03,
        1.07,
        f"注意｜{trap}",
        fill=WARNING_BG,
        line=lighten(WARNING, 0.55),
        size=12.5,
        color=WARNING,
        bold=False,
    )


def render_cards(slide, items, accent, soft) -> None:
    items = items[:6]
    cols = 3 if len(items) >= 5 else 2
    rows = (len(items) + cols - 1) // cols
    gap_x, gap_y = 0.22, 0.22
    total_w, total_h = 12.25, 4.42
    card_w = (total_w - gap_x * (cols - 1)) / cols
    card_h = (total_h - gap_y * (rows - 1)) / rows
    for index, (label, body) in enumerate(items):
        row, col = divmod(index, cols)
        add_card(
            slide,
            0.55 + col * (card_w + gap_x),
            1.12 + row * (card_h + gap_y),
            card_w,
            card_h,
            label,
            body,
            accent,
            soft,
        )


def render_process(slide, items, accent, soft) -> None:
    items = items[:5]
    count = len(items)
    arrow_w = 0.34
    gap = 0.12
    total_arrow_w = arrow_w * max(0, count - 1)
    width = (12.05 - total_arrow_w - gap * 2 * max(0, count - 1)) / max(count, 1)
    for index, (label, body) in enumerate(items):
        x = 0.62 + index * (width + arrow_w + gap * 2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.58),
            Inches(width),
            Inches(3.35),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(soft if index % 2 == 0 else WHITE)
        shape.line.color.rgb = rgb(accent)
        shape.line.width = Pt(1.4)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + width / 2 - 0.28),
            Inches(1.34),
            Inches(0.56),
            Inches(0.56),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(accent)
        badge.line.fill.background()
        set_shape_text(
            badge,
            str(index + 1),
            size=13,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + 0.13,
            1.98,
            width - 0.26,
            0.74,
            shorten(label, 16),
            size=13.5,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + 0.16,
            2.78,
            width - 0.32,
            1.70,
            shorten(body, 42),
            size=11.5,
            color=INK,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.TOP,
        )
        if index < count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + width + gap),
                Inches(2.84),
                Inches(arrow_w),
                Inches(0.70),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(lighten(accent, 0.28))
            arrow.line.fill.background()


def render_comparison(slide, items, accent, soft) -> None:
    midpoint = (len(items) + 1) // 2
    columns = [items[:midpoint], items[midpoint:]]
    headings = ["視点 A", "視点 B"]
    for col, group in enumerate(columns):
        x = 0.62 + col * 6.18
        add_box(
            slide,
            x,
            1.15,
            5.90,
            0.46,
            headings[col],
            fill=accent if col == 0 else soft,
            line=accent,
            size=14,
            color=WHITE if col == 0 else accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        card_h = min(1.18, 3.72 / max(len(group), 1))
        for row, (label, body) in enumerate(group[:3]):
            add_card(
                slide,
                x,
                1.76 + row * (card_h + 0.12),
                5.90,
                card_h,
                label,
                body,
                accent,
                soft,
            )


def render_matrix(slide, items, accent, soft) -> None:
    values = (items + [("確認", "本文で条件を整理")] * 4)[:4]
    positions = [(0.80, 1.25), (6.82, 1.25), (0.80, 3.35), (6.82, 3.35)]
    for index, ((label, body), (x, y)) in enumerate(zip(values, positions)):
        add_card(slide, x, y, 5.72, 1.78, label, body, accent, soft)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.20), Inches(y - 0.14), Inches(0.48), Inches(0.48)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(accent)
        badge.line.fill.background()
        set_shape_text(
            badge,
            str(index + 1),
            size=12,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def render_cycle(slide, items, accent, soft) -> None:
    values = (items + [("確認", "次の段階へ")] * 4)[:4]
    positions = [(4.95, 1.15), (8.72, 2.50), (4.95, 4.15), (1.18, 2.50)]
    center = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.35), Inches(2.52), Inches(2.62), Inches(1.55)
    )
    center.fill.solid()
    center.fill.fore_color.rgb = rgb(accent)
    center.line.fill.background()
    set_shape_text(
        center,
        "循環して\n改善・定着",
        size=17,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    for index, ((label, body), (x, y)) in enumerate(zip(values, positions)):
        add_card(slide, x, y, 3.18, 1.12, label, body, accent, soft)
        next_x, next_y = positions[(index + 1) % 4]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x + 1.59),
            Inches(y + 0.56),
            Inches(next_x + 1.59),
            Inches(next_y + 0.56),
        )
        connector.line.color.rgb = rgb(lighten(accent, 0.35))
        connector.line.width = Pt(1.8)


def render_hierarchy(slide, items, accent, soft) -> None:
    values = items[:5]
    for index, (label, body) in enumerate(values):
        width = 11.5 - index * 1.25
        x = (13.333 - width) / 2
        y = 1.12 + index * 0.88
        add_box(
            slide,
            x,
            y,
            width,
            0.70,
            f"{index + 1}｜{label}　{shorten(body, 55)}",
            fill=soft if index % 2 == 0 else WHITE,
            line=accent,
            size=13.5,
            color=INK,
            bold=index == 0,
            align=PP_ALIGN.CENTER,
        )


def render_pyramid(slide, items, accent, soft) -> None:
    values = items[:5]
    values = list(reversed(values))
    for index, (label, body) in enumerate(values):
        width = 3.5 + index * 1.75
        x = (13.333 - width) / 2
        y = 1.10 + index * 0.88
        shape = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID, Inches(x), Inches(y), Inches(width), Inches(0.74)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(soft if index % 2 == 0 else WHITE)
        shape.line.color.rgb = rgb(accent)
        set_shape_text(
            shape,
            f"{shorten(label, 18)}｜{shorten(body, 38)}",
            size=12.5,
            color=INK,
            align=PP_ALIGN.CENTER,
        )


def render_timeline(slide, items, accent, soft) -> None:
    values = items[:5]
    count = len(values)
    y = 3.20
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y), Inches(11.45), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(accent)
    line.line.fill.background()
    for index, (label, body) in enumerate(values):
        x = 0.85 + index * (11.25 / max(count - 1, 1))
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y - 0.18), Inches(0.42), Inches(0.42)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(accent)
        dot.line.fill.background()
        box_y = 1.40 if index % 2 == 0 else 3.72
        add_box(
            slide,
            max(0.42, min(x - 0.90, 10.92)),
            box_y,
            2.25,
            1.18,
            f"{index + 1}｜{shorten(label, 18)}\n{shorten(body, 38)}",
            fill=soft if index % 2 == 0 else WHITE,
            line=accent,
            size=12,
            color=INK,
            align=PP_ALIGN.CENTER,
        )


def render_network(slide, items, accent, soft, center_text: str) -> None:
    values = items[:6]
    center_x, center_y = 5.15, 2.18
    center = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(center_x), Inches(center_y), Inches(3.05), Inches(1.65)
    )
    center.fill.solid()
    center.fill.fore_color.rgb = rgb(accent)
    center.line.fill.background()
    set_shape_text(
        center,
        shorten(center_text, 24),
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    positions = [
        (0.60, 1.15),
        (9.72, 1.15),
        (0.60, 3.90),
        (9.72, 3.90),
        (3.05, 4.48),
        (7.25, 4.48),
    ]
    for (label, body), (x, y) in zip(values, positions):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(center_x + 1.52),
            Inches(center_y + 0.82),
            Inches(x + 1.50),
            Inches(y + 0.50),
        )
        connector.line.color.rgb = rgb(lighten(accent, 0.35))
        connector.line.width = Pt(1.5)
        add_card(slide, x, y, 3.0, 1.0, label, body, accent, soft)


def render_equation(slide, items, accent, soft, equation: str) -> None:
    add_box(
        slide,
        0.85,
        1.30,
        11.63,
        1.28,
        equation or "定義・関係式を確認",
        fill=soft,
        line=accent,
        size=22,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    values = items[:3]
    width = 3.72
    for index, (label, body) in enumerate(values):
        add_card(
            slide,
            0.85 + index * 3.93,
            2.92,
            width,
            2.20,
            label,
            body,
            accent,
            soft,
        )


def render_table(slide, topic: Topic, items, accent, soft) -> None:
    rows = markdown_table(topic.sections["core"])
    if len(rows) < 2:
        render_cards(slide, items, accent, soft)
        return
    headers = rows[0][:4]
    data = rows[1:7]
    col_count = max(2, len(headers))
    table_shape = slide.shapes.add_table(
        len(data) + 1,
        col_count,
        Inches(0.65),
        Inches(1.22),
        Inches(12.03),
        Inches(4.30),
    )
    table = table_shape.table
    for col in range(col_count):
        table.columns[col].width = Inches(12.03 / col_count)
    for row_index in range(len(data) + 1):
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(
                accent if row_index == 0 else (soft if row_index % 2 else WHITE)
            )
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            value = ""
            if row_index == 0 and col_index < len(headers):
                value = headers[col_index]
            elif row_index > 0 and col_index < len(data[row_index - 1]):
                value = data[row_index - 1][col_index]
            set_shape_text(
                cell,
                shorten(value, 44),
                size=11.5,
                color=WHITE if row_index == 0 else INK,
                bold=row_index == 0 or col_index == 0,
                align=PP_ALIGN.CENTER,
                margin=0.06,
            )


def render_visual(slide, topic: Topic, accent: str, soft: str) -> str:
    items = key_items(topic)
    visual = diagram_type(topic)
    renderers = {
        "cards": lambda: render_cards(slide, items, accent, soft),
        "process": lambda: render_process(slide, items, accent, soft),
        "comparison": lambda: render_comparison(slide, items, accent, soft),
        "matrix": lambda: render_matrix(slide, items, accent, soft),
        "cycle": lambda: render_cycle(slide, items, accent, soft),
        "hierarchy": lambda: render_hierarchy(slide, items, accent, soft),
        "pyramid": lambda: render_pyramid(slide, items, accent, soft),
        "timeline": lambda: render_timeline(slide, items, accent, soft),
        "network": lambda: render_network(
            slide, items, accent, soft, topic.title.split("（")[0]
        ),
        "equation": lambda: render_equation(
            slide, items, accent, soft, formula_text(topic)
        ),
        "table": lambda: render_table(slide, topic, items, accent, soft),
    }
    renderers[visual]()
    return visual


def add_title_slide(prs: Presentation, book: Textbook, config: dict) -> None:
    accent, soft = config["accent"], config["soft"]
    slide = add_slide_base(prs, accent)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.65), Inches(0.75), Inches(4.20), Inches(4.20)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(soft)
    circle.line.color.rgb = rgb(lighten(accent, 0.50))
    circle.line.width = Pt(2)
    set_shape_text(
        circle,
        book.subject,
        size=54,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        0.75,
        1.40,
        7.75,
        0.55,
        "中小企業診断士 第1次試験",
        size=17,
        color=accent,
        bold=True,
    )
    add_text(
        slide,
        0.72,
        2.05,
        8.10,
        1.35,
        config["name"],
        size=36,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        0.75,
        3.65,
        7.4,
        0.70,
        "図解スライド｜全論点を「構造・比較・流れ」で理解する",
        size=20,
        color=MUTED,
    )
    add_box(
        slide,
        0.75,
        5.05,
        3.55,
        0.72,
        f"{len(book.topics)} 論点",
        fill=accent,
        line=accent,
        size=20,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_box(
        slide,
        4.48,
        5.05,
        3.55,
        0.72,
        f"{len(book.parts)} セクション",
        fill=soft,
        line=accent,
        size=20,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        0.75,
        6.43,
        11.9,
        0.28,
        "SMEC 教科書シリーズ｜教科書本文を基に自動生成",
        size=11,
        color=MUTED,
    )


def add_howto_slide(prs: Presentation, book: Textbook, config: dict) -> None:
    accent, soft = config["accent"], config["soft"]
    slide = add_slide_base(prs, accent)
    add_text(
        slide,
        0.65,
        0.42,
        12.0,
        0.55,
        "このスライドの使い方",
        size=28,
        color=INK,
        bold=True,
    )
    steps = [
        ("1", "全体像", "まず図の関係性を眺める"),
        ("2", "具体例", "身近な場面へ置き換える"),
        ("3", "注意", "誤答パターンを見分ける"),
        ("4", "演習", "教科書の参照先で定着させる"),
    ]
    for index, (number, label, body) in enumerate(steps):
        x = 0.72 + index * 3.12
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.92), Inches(1.38), Inches(0.75), Inches(0.75)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(accent)
        badge.line.fill.background()
        set_shape_text(
            badge,
            number,
            size=20,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_card(slide, x, 2.32, 2.60, 1.72, label, body, accent, soft)
        if index < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + 2.67),
                Inches(2.86),
                Inches(0.38),
                Inches(0.55),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(lighten(accent, 0.30))
            arrow.line.fill.background()
    intro = section_text(book.intro, 180)
    add_box(
        slide,
        0.75,
        4.65,
        11.83,
        1.40,
        intro or "各論点を図で理解した後、教科書本文と過去問で知識を定着させます。",
        fill=soft,
        line=lighten(accent, 0.45),
        size=15,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        0.78,
        6.30,
        11.7,
        0.35,
        "※ 数値・法令・制度は出題年度によって変わる場合があります。最新の公式情報も確認してください。",
        size=11,
        color=WARNING,
        align=PP_ALIGN.CENTER,
    )


def add_part_slide(
    prs: Presentation, part: str, index: int, topic_count: int, config: dict
) -> None:
    accent, soft = config["accent"], config["soft"]
    slide = add_slide_base(prs, accent)
    block = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(0.92),
        Inches(11.63),
        Inches(5.55),
    )
    block.fill.solid()
    block.fill.fore_color.rgb = rgb(soft)
    block.line.color.rgb = rgb(lighten(accent, 0.45))
    block.line.width = Pt(2)
    add_text(
        slide,
        1.35,
        1.50,
        2.2,
        0.42,
        f"SECTION {index:02d}",
        size=15,
        color=accent,
        bold=True,
    )
    add_text(
        slide,
        1.30,
        2.18,
        10.70,
        1.45,
        re.sub(r"（[A-G]-\d+〜[A-G]-\d+）", "", part),
        size=33,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_box(
        slide,
        4.66,
        4.34,
        4.0,
        0.82,
        f"{topic_count} 論点",
        fill=accent,
        line=accent,
        size=21,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_summary_slide(prs: Presentation, book: Textbook, config: dict) -> None:
    accent, soft = config["accent"], config["soft"]
    slide = add_slide_base(prs, accent)
    add_text(
        slide,
        0.65,
        0.42,
        12.0,
        0.55,
        f"{config['name']}｜学習のまとめ",
        size=27,
        color=INK,
        bold=True,
    )
    summary_items: list[tuple[str, str]] = []
    for line in book.summary:
        stripped = line.strip()
        if re.match(r"^(?:[-*+]|\d+[.)])\s+", stripped):
            summary_items.append(split_label_body(stripped))
    if not summary_items:
        summary_items = [
            ("理解", "図で概念同士の関係を捉える"),
            ("確認", "教科書本文で定義と条件を確認する"),
            ("演習", "問題を解いて知識を使える形にする"),
            ("復習", "注意ポイントから弱点へ戻る"),
        ]
    render_cycle(slide, summary_items[:4], accent, soft)
    add_box(
        slide,
        0.85,
        5.55,
        11.65,
        1.06,
        f"次の一歩｜problem_sets/1st_stage/ の論点表と exercises/ の演習で、{len(book.topics)}論点を定着させる",
        fill=soft,
        line=accent,
        size=14.5,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_deck(book: Textbook, output: Path) -> dict[str, int]:
    config = SUBJECTS[book.subject]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = f"{config['name']} 図解スライド"
    prs.core_properties.subject = "中小企業診断士 第1次試験"
    prs.core_properties.author = "SMEC Project"
    prs.core_properties.keywords = "中小企業診断士, 教科書, 図解, スライド"

    add_title_slide(prs, book, config)
    add_howto_slide(prs, book, config)
    slide_number = 2
    diagram_counts: dict[str, int] = {}

    grouped: dict[str, list[Topic]] = {part: [] for part in book.parts}
    for topic in book.topics:
        grouped.setdefault(topic.part, []).append(topic)

    for part_index, (part, topics) in enumerate(grouped.items(), start=1):
        if not topics:
            continue
        add_part_slide(prs, part, part_index, len(topics), config)
        slide_number += 1
        for topic in topics:
            slide = add_slide_base(prs, config["accent"])
            add_topic_header(slide, topic, config["accent"])
            visual = render_visual(slide, topic, config["accent"], config["soft"])
            diagram_counts[visual] = diagram_counts.get(visual, 0) + 1
            add_context_panels(slide, topic, config["accent"])
            slide_number += 1
            add_footer(slide, config["file"], slide_number, config["accent"])

    add_summary_slide(prs, book, config)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return diagram_counts


def write_readme(books: list[Textbook], outputs: list[Path]) -> None:
    rows = [
        "# 教科書 図解スライド",
        "",
        "各教科書の全論点を、比較・フロー・マトリクス・階層・循環などの図解で整理したPowerPoint資料です。",
        "",
        "| 科目 | 元の教科書 | スライド | 論点数 |",
        "|---|---|---|---:|",
    ]
    for book, output in zip(books, outputs):
        config = SUBJECTS[book.subject]
        rows.append(
            f"| {book.subject}. {config['name']} | "
            f"[`{config['file']}`](../{config['file']}) | "
            f"[`{output.name}`]({output.name}) | {len(book.topics)} |"
        )
    rows.extend(
        [
            "",
            "## 再生成",
            "",
            "```bash",
            "python3 -m pip install -r scripts/requirements-slides.txt",
            "python3 scripts/generate_textbook_slides.py",
            "```",
            "",
            "特定科目だけ生成する場合は、科目記号を指定します。",
            "",
            "```bash",
            "python3 scripts/generate_textbook_slides.py --subject C",
            "```",
            "",
            "スライドは教科書本文を基に生成しています。法令・統計・制度の数値は、受験年度の公式情報も確認してください。",
            "",
        ]
    )
    (OUTPUT_DIR / "README.md").write_text("\n".join(rows), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--subject",
        choices=sorted(SUBJECTS),
        action="append",
        help="生成する科目記号。複数回指定可能（省略時はA科目）",
    )
    parser.add_argument(
        "--all",
        action="store_true",
        help="全7科目を生成する",
    )
    args = parser.parse_args()
    if args.all and args.subject:
        parser.error("--all と --subject は同時に指定できません")
    subjects = list(SUBJECTS) if args.all else (args.subject or ["A"])

    for subject in subjects:
        config = SUBJECTS[subject]
        source = TEXTBOOK_DIR / config["file"]
        book = parse_textbook(source, subject)
        output = OUTPUT_DIR / source.name.replace("_textbook.md", "_visual_slides.pptx")
        counts = build_deck(book, output)
        visual_summary = ", ".join(f"{key}:{value}" for key, value in sorted(counts.items()))
        print(f"{subject}: {len(book.topics)}論点 → {output} ({visual_summary})")

    existing_books: list[Textbook] = []
    existing_outputs: list[Path] = []
    for subject, config in SUBJECTS.items():
        output = OUTPUT_DIR / config["file"].replace("_textbook.md", "_visual_slides.pptx")
        if output.exists():
            existing_books.append(
                parse_textbook(TEXTBOOK_DIR / config["file"], subject)
            )
            existing_outputs.append(output)
    write_readme(existing_books, existing_outputs)


if __name__ == "__main__":
    main()
