#!/usr/bin/env python3
"""SMEC向けの汎用図解12種＋経済学専用グラフ6種を生成する。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "textbook" / "slides" / "SMEC_visual_diagram_template.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "WenQuanYi Micro Hei"

INK = "161616"
NAVY = "0F172A"
BLUE = "0F62FE"
CYAN = "1192E8"
TEAL = "009D9A"
GREEN = "24A148"
YELLOW = "F1C21B"
ORANGE = "FF832B"
RED = "DA1E28"
PURPLE = "8A3FFC"
GRAY_10 = "F4F4F4"
GRAY_20 = "E0E0E0"
GRAY_50 = "8D8D8D"
GRAY_70 = "525252"
WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 16,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.04,
):
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str = "",
    *,
    fill: str = WHITE,
    line: str = GRAY_20,
    size: float = 15,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    radius: bool = True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    if text:
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.10)
        frame.margin_right = Inches(0.10)
        frame.margin_top = Inches(0.07)
        frame.margin_bottom = Inches(0.07)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = GRAY_50,
    width: float = 1.5,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    return connector


def add_circle(
    slide,
    x: float,
    y: float,
    d: float,
    text: str,
    *,
    fill: str = BLUE,
    color: str = WHITE,
    size: float = 15,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)
    return shape


def add_base(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    return slide


def add_header(
    slide,
    code: str,
    category: str,
    action_title: str,
    *,
    accent: str = BLUE,
) -> None:
    add_box(
        slide,
        0.50,
        0.30,
        0.92,
        0.36,
        code,
        fill=accent,
        line=accent,
        size=12,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        1.60,
        0.30,
        2.40,
        0.36,
        category,
        size=11,
        color=GRAY_70,
        bold=True,
    )
    add_text(
        slide,
        0.50,
        0.78,
        12.15,
        0.62,
        action_title,
        size=24,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_line(slide, 0.50, 1.48, 12.83, 1.48, color=GRAY_20, width=1)


def add_takeaway(
    slide,
    text: str,
    number: int,
    *,
    accent: str = BLUE,
) -> None:
    add_box(
        slide,
        0.50,
        6.50,
        12.33,
        0.55,
        f"KEY TAKEAWAY｜{text}",
        fill=GRAY_10,
        line=GRAY_20,
        size=11.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.LEFT,
        radius=False,
    )
    add_text(
        slide,
        0.52,
        7.16,
        10.5,
        0.18,
        "SMEC Visual Slide System｜すべて編集可能なPowerPoint図形",
        size=8,
        color=GRAY_50,
    )
    add_text(
        slide,
        12.15,
        7.14,
        0.60,
        0.20,
        f"{number:02d}",
        size=9,
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_title_slide(prs: Presentation) -> None:
    slide = add_base(prs)
    add_box(
        slide,
        0,
        0,
        4.10,
        7.50,
        "",
        fill=NAVY,
        line=NAVY,
        radius=False,
    )
    for x, y, d, fill in [
        (0.60, 0.80, 1.10, BLUE),
        (2.10, 1.40, 0.68, CYAN),
        (1.15, 2.65, 1.50, TEAL),
        (2.70, 3.85, 0.95, PURPLE),
        (0.60, 5.15, 1.25, ORANGE),
    ]:
        add_circle(slide, x, y, d, "", fill=fill)
    add_text(
        slide,
        4.85,
        1.05,
        7.55,
        0.35,
        "SMEC DESIGN SYSTEM",
        size=15,
        color=BLUE,
        bold=True,
    )
    add_text(
        slide,
        4.80,
        1.65,
        7.70,
        1.45,
        "図で理解するための\n標準スライドテンプレート",
        size=34,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        4.85,
        3.45,
        7.20,
        0.80,
        "汎用図解 12パターン\n＋ 経済学専用グラフ 6パターン",
        size=21,
        color=GRAY_70,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_box(
        slide,
        4.85,
        5.25,
        2.05,
        0.52,
        "18 PATTERNS",
        fill=BLUE,
        line=BLUE,
        size=13,
        color=WHITE,
        bold=True,
    )
    add_box(
        slide,
        7.10,
        5.25,
        2.05,
        0.52,
        "100% EDITABLE",
        fill=WHITE,
        line=BLUE,
        size=13,
        color=BLUE,
        bold=True,
    )
    add_text(
        slide,
        4.85,
        6.65,
        7.30,
        0.28,
        "公開デザインシステムを参考にしたSMEC独自仕様",
        size=11,
        color=GRAY_50,
    )


def add_design_rules(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "RULES",
        "デザイン原則",
        "統一感は「グリッド・色・文字・1枚1メッセージ」の固定から生まれる",
    )
    cards = [
        ("01", "ACTION TITLE", "タイトルだけで結論が伝わる\n名詞ではなく文章で書く", BLUE),
        ("02", "ONE MESSAGE", "1枚につき主張は1つ\n図は主張を証明するために置く", TEAL),
        ("03", "GRID SYSTEM", "12列グリッドを基準に整列\n余白は0.5インチ以上", PURPLE),
        ("04", "COLOR ROLE", "青＝主張／緑＝望ましい\n赤＝注意／灰＝補助情報", ORANGE),
    ]
    for i, (num, title, body, color) in enumerate(cards):
        x = 0.58 + i * 3.12
        add_box(
            slide,
            x,
            1.82,
            2.85,
            3.72,
            "",
            fill=WHITE,
            line=GRAY_20,
        )
        add_circle(slide, x + 0.22, 2.08, 0.62, num, fill=color, size=13)
        add_text(
            slide,
            x + 0.22,
            2.92,
            2.40,
            0.36,
            title,
            size=14,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            x + 0.22,
            3.48,
            2.40,
            1.26,
            body,
            size=14,
            color=INK,
            valign=MSO_ANCHOR.TOP,
        )
    add_takeaway(
        slide,
        "装飾を増やすより、同じルールを全ページで守る方がプロ品質に近づく。",
        len(prs.slides),
    )


def add_selector(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "INDEX",
        "パターン選択",
        "伝えたい関係を先に決めれば、使う図解パターンは自動的に絞り込める",
    )
    rows = [
        ("構成要素を並べる", "G01 定義カード / G09 ツリー"),
        ("違いを示す", "G02 比較 / G08 2×2マトリクス"),
        ("順序や変化を示す", "G03 プロセス / G07 タイムライン"),
        ("理由や影響を示す", "G04 因果 / G05 循環"),
        ("レベルや関係を示す", "G06 階層 / G10 ネットワーク"),
        ("数値・式で示す", "G11 数式 / G12 KPI・表"),
        ("経済グラフで示す", "E01〜E06 専用パターン"),
    ]
    for i, (question, answer) in enumerate(rows):
        y = 1.76 + i * 0.62
        add_box(
            slide,
            0.68,
            y,
            4.10,
            0.46,
            question,
            fill=GRAY_10,
            line=GRAY_20,
            size=12.5,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.LEFT,
            radius=False,
        )
        add_box(
            slide,
            4.95,
            y,
            7.70,
            0.46,
            answer,
            fill=WHITE,
            line=GRAY_20,
            size=12.5,
            color=BLUE,
            bold=True,
            align=PP_ALIGN.LEFT,
            radius=False,
        )
    add_takeaway(
        slide,
        "内容に図を当てはめるのではなく、「何の関係を説明するか」から型を選ぶ。",
        len(prs.slides),
    )


def pattern_definition(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G01",
        "定義・要素カード",
        "複雑な概念も「定義・構成・判断基準」の3点に分けると理解しやすい",
    )
    items = [
        ("定義", "GDPは国内で新たに生み出された付加価値の合計", BLUE),
        ("構成", "生産・分配・支出の3面から見ても合計は等しい", TEAL),
        ("判断基準", "新しい価値を生んだか、市場取引されたかを確認", PURPLE),
    ]
    for i, (label, body, color) in enumerate(items):
        x = 0.68 + i * 4.12
        add_box(slide, x, 1.88, 3.82, 3.85, "", fill=WHITE, line=GRAY_20)
        add_circle(slide, x + 1.51, 2.16, 0.80, str(i + 1), fill=color, size=18)
        add_text(
            slide,
            x + 0.28,
            3.18,
            3.26,
            0.46,
            label,
            size=18,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + 0.35,
            3.92,
            3.12,
            1.05,
            body,
            size=14,
            color=INK,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.TOP,
        )
    add_takeaway(
        slide,
        "新概念の導入、用語説明、章の冒頭に使う。カード数は3〜5個まで。",
        len(prs.slides),
    )


def pattern_comparison(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G02",
        "左右比較",
        "似た概念は共通の比較軸に揃えると、違いを一目で判別できる",
    )
    add_box(
        slide,
        0.72,
        1.78,
        5.70,
        0.62,
        "名目GDP",
        fill=BLUE,
        line=BLUE,
        size=20,
        color=WHITE,
        bold=True,
    )
    add_box(
        slide,
        6.92,
        1.78,
        5.70,
        0.62,
        "実質GDP",
        fill=TEAL,
        line=TEAL,
        size=20,
        color=WHITE,
        bold=True,
    )
    comparisons = [
        ("価格", "その年の価格", "基準年の価格"),
        ("物価変動", "含む", "取り除く"),
        ("主な用途", "経済規模の金額比較", "実質的な成長率の比較"),
    ]
    for i, (axis, left, right) in enumerate(comparisons):
        y = 2.68 + i * 0.88
        add_box(
            slide,
            0.72,
            y,
            1.40,
            0.62,
            axis,
            fill=GRAY_10,
            line=GRAY_20,
            size=12,
            color=GRAY_70,
            bold=True,
            radius=False,
        )
        add_box(
            slide,
            2.20,
            y,
            4.22,
            0.62,
            left,
            fill=WHITE,
            line=GRAY_20,
            size=14,
            color=NAVY,
            radius=False,
        )
        add_box(
            slide,
            6.92,
            y,
            5.70,
            0.62,
            right,
            fill=WHITE,
            line=GRAY_20,
            size=14,
            color=NAVY,
            radius=False,
        )
    add_takeaway(
        slide,
        "「AとBの違い」は、両側で同じ比較軸を使い、対称配置にする。",
        len(prs.slides),
    )


def pattern_process(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G03",
        "プロセス・手順",
        "判断手順は左から右へ流し、各段階の入力と結果を明示する",
    )
    steps = [
        ("状況確認", "数値・条件を読む"),
        ("モデル選択", "使う理論を決める"),
        ("方向判定", "曲線のシフトを見る"),
        ("結論", "価格・数量を答える"),
    ]
    for i, (label, body) in enumerate(steps):
        x = 0.62 + i * 3.08
        add_circle(slide, x + 0.90, 2.04, 0.62, str(i + 1), fill=BLUE, size=14)
        add_box(
            slide,
            x,
            2.90,
            2.55,
            2.10,
            "",
            fill=GRAY_10 if i % 2 == 0 else WHITE,
            line=BLUE,
        )
        add_text(
            slide,
            x + 0.16,
            3.22,
            2.23,
            0.38,
            label,
            size=16,
            color=BLUE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + 0.25,
            3.90,
            2.05,
            0.55,
            body,
            size=13,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + 2.65),
                Inches(3.58),
                Inches(0.34),
                Inches(0.70),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(CYAN)
            arrow.line.fill.background()
    add_takeaway(
        slide,
        "手続、解法、業務フローに使う。4〜6段階で、動詞から書き始める。",
        len(prs.slides),
    )


def pattern_cause_effect(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G04",
        "因果関係",
        "原因から結果までの連鎖を分解すると、経済政策の波及経路が見える",
    )
    nodes = [
        (0.70, 2.12, 2.35, 0.90, "原因\n政府支出が増加", BLUE),
        (3.55, 2.12, 2.35, 0.90, "一次効果\n総需要が増加", CYAN),
        (6.40, 2.12, 2.35, 0.90, "波及効果\n所得・消費が増加", TEAL),
        (9.25, 2.12, 3.15, 0.90, "結果\nGDPが乗数倍増加", GREEN),
    ]
    for i, (x, y, w, h, text, color) in enumerate(nodes):
        add_box(
            slide,
            x,
            y,
            w,
            h,
            text,
            fill=color,
            line=color,
            size=14,
            color=WHITE,
            bold=True,
        )
        if i < len(nodes) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + w + 0.12),
                Inches(y + 0.18),
                Inches(0.45),
                Inches(0.54),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(GRAY_50)
            arrow.line.fill.background()
    add_box(
        slide,
        3.10,
        4.10,
        7.15,
        1.08,
        "反作用｜利子率上昇 → 民間投資減少（クラウディングアウト）",
        fill="FFF1F1",
        line=RED,
        size=16,
        color=RED,
        bold=True,
    )
    add_takeaway(
        slide,
        "原因・媒介・結果・反作用を分けると、暗記ではなく仕組みで理解できる。",
        len(prs.slides),
    )


def pattern_cycle(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G05",
        "循環モデル",
        "繰り返される活動は、中心概念の周囲に4段階で配置すると記憶に残る",
    )
    center_x, center_y = 5.32, 2.70
    add_circle(slide, center_x, center_y, 2.10, "景気循環", fill=NAVY, size=21)
    cycle_nodes = [
        (5.53, 1.66, "回復", BLUE),
        (8.58, 2.96, "好況", GREEN),
        (5.53, 4.76, "後退", ORANGE),
        (2.45, 2.96, "不況", PURPLE),
    ]
    for i, (x, y, label, color) in enumerate(cycle_nodes):
        add_box(
            slide,
            x,
            y,
            1.68,
            0.72,
            label,
            fill=color,
            line=color,
            size=17,
            color=WHITE,
            bold=True,
        )
        next_x, next_y, _, _ = cycle_nodes[(i + 1) % len(cycle_nodes)]
        add_line(
            slide,
            x + 0.84,
            y + 0.36,
            next_x + 0.84,
            next_y + 0.36,
            color=GRAY_50,
            width=1.8,
        )
    add_takeaway(
        slide,
        "PDCA、景気循環、フィードバックに使う。時計回りに統一する。",
        len(prs.slides),
    )


def pattern_hierarchy(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G06",
        "階層・ピラミッド",
        "上位概念から具体例へ絞り込むと、包含関係と優先順位を同時に示せる",
    )
    levels = [
        (3.98, 1.88, 5.38, 0.72, "経済学", NAVY, WHITE),
        (2.90, 2.83, 7.55, 0.72, "マクロ経済学 ｜ ミクロ経済学", BLUE, WHITE),
        (1.80, 3.78, 9.75, 0.72, "市場・政策・国際・消費者・企業・分配", CYAN, WHITE),
        (0.70, 4.73, 11.95, 0.72, "36の個別論点と過去問パターン", GRAY_10, NAVY),
    ]
    for x, y, w, h, text, fill, color in levels:
        add_box(
            slide,
            x,
            y,
            w,
            h,
            text,
            fill=fill,
            line=fill if fill != GRAY_10 else GRAY_20,
            size=16,
            color=color,
            bold=True,
            radius=False,
        )
    add_takeaway(
        slide,
        "概念分類、戦略レベル、優先順位に使う。上位ほど抽象、下位ほど具体。",
        len(prs.slides),
    )


def pattern_timeline(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G07",
        "タイムライン",
        "出来事を時間軸に置くと、理論が生まれた背景と変化の順序を捉えられる",
    )
    add_line(slide, 1.05, 3.56, 12.25, 3.56, color=BLUE, width=3)
    events = [
        ("1930s", "世界恐慌", "ケインズ理論", BLUE),
        ("1970s", "石油危機", "スタグフレーション", ORANGE),
        ("1980s", "政策転換", "マネタリズム", PURPLE),
        ("2000s", "低金利", "非伝統的金融政策", TEAL),
    ]
    for i, (year, event, theory, color) in enumerate(events):
        x = 0.98 + i * 3.08
        add_circle(slide, x + 0.90, 3.28, 0.58, "", fill=color)
        box_y = 1.85 if i % 2 == 0 else 4.12
        add_box(
            slide,
            x,
            box_y,
            2.38,
            1.30,
            "",
            fill=WHITE,
            line=color,
        )
        add_text(
            slide,
            x + 0.16,
            box_y + 0.13,
            2.06,
            0.30,
            year,
            size=14,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + 0.15,
            box_y + 0.54,
            2.08,
            0.58,
            f"{event}\n{theory}",
            size=12,
            color=INK,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.TOP,
        )
    add_takeaway(
        slide,
        "歴史、制度変更、ライフサイクルに使う。時間は必ず左から右へ流す。",
        len(prs.slides),
    )


def pattern_matrix(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G08",
        "2×2マトリクス",
        "2つの判断軸を交差させると、4つの状態を漏れなく整理できる",
    )
    x, y, w, h = 2.32, 1.78, 8.72, 4.28
    add_box(slide, x, y, w / 2, h / 2, "高成長 × 低失業\n需要超過・インフレ注意", fill="E8F1FF", line=WHITE, size=16, color=BLUE, bold=True, radius=False)
    add_box(slide, x + w / 2, y, w / 2, h / 2, "高成長 × 高失業\n構造的失業を点検", fill="E5F6FF", line=WHITE, size=16, color=CYAN, bold=True, radius=False)
    add_box(slide, x, y + h / 2, w / 2, h / 2, "低成長 × 低失業\n供給制約を点検", fill="D9FBFB", line=WHITE, size=16, color=TEAL, bold=True, radius=False)
    add_box(slide, x + w / 2, y + h / 2, w / 2, h / 2, "低成長 × 高失業\n景気刺激策を検討", fill="FFF1F1", line=WHITE, size=16, color=RED, bold=True, radius=False)
    add_text(slide, 0.76, 3.17, 1.28, 0.42, "成長率", size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 5.55, 6.08, 2.20, 0.30, "失業率 →", size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(
        slide,
        "軸は独立した2変数にする。各象限の名称だけでなく判断も添える。",
        len(prs.slides),
    )


def pattern_tree(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G09",
        "ツリー・分類",
        "概念をMECEに枝分かれさせると、全体像と用語の所属を同時に確認できる",
    )
    add_box(slide, 0.72, 3.05, 2.16, 0.82, "物価指数", fill=NAVY, line=NAVY, size=18, color=WHITE, bold=True)
    categories = [
        (4.10, 1.78, "CPI", "消費者が購入する財・サービス", BLUE),
        (4.10, 3.05, "CGPI", "企業間で取引される商品", TEAL),
        (4.10, 4.32, "GDPデフレーター", "GDP全体の物価動向", PURPLE),
    ]
    for x, y, title, body, color in categories:
        add_line(slide, 2.88, 3.46, x, y + 0.44, color=GRAY_50, width=1.5)
        add_box(slide, x, y, 2.32, 0.88, title, fill=color, line=color, size=15, color=WHITE, bold=True)
        add_box(slide, 6.75, y, 5.35, 0.88, body, fill=GRAY_10, line=GRAY_20, size=13, color=NAVY, align=PP_ALIGN.LEFT)
        add_line(slide, 6.42, y + 0.44, 6.75, y + 0.44, color=GRAY_50)
    add_takeaway(
        slide,
        "分類、論点体系、原因分解に使う。同じ階層では分類基準を混ぜない。",
        len(prs.slides),
    )


def pattern_network(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G10",
        "ネットワーク・関係図",
        "中心概念と周辺主体の相互作用を示すと、複雑な市場関係を俯瞰できる",
    )
    add_circle(slide, 5.37, 2.58, 2.10, "市場", fill=NAVY, size=22)
    nodes = [
        (0.85, 1.82, "消費者", BLUE),
        (9.78, 1.82, "企業", TEAL),
        (0.85, 4.35, "政府", PURPLE),
        (9.78, 4.35, "海外", ORANGE),
    ]
    for x, y, label, color in nodes:
        add_line(slide, 6.42, 3.63, x + 1.35, y + 0.55, color=GRAY_50, width=1.6)
        add_box(slide, x, y, 2.70, 1.10, label, fill=color, line=color, size=18, color=WHITE, bold=True)
    add_box(slide, 4.02, 5.18, 4.80, 0.64, "価格・数量・所得・情報が相互に流れる", fill=GRAY_10, line=GRAY_20, size=13, color=NAVY, bold=True)
    add_takeaway(
        slide,
        "ステークホルダー、5フォース、SCMに使う。線の意味は凡例で明示する。",
        len(prs.slides),
    )


def pattern_formula(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G11",
        "数式・投入産出",
        "数式は「入力→関係式→出力」に分けると、記号の意味と計算手順が見える",
    )
    add_box(slide, 0.76, 2.15, 3.10, 2.70, "", fill=GRAY_10, line=GRAY_20)
    add_text(slide, 1.02, 2.40, 2.58, 0.36, "INPUT", size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 3.08, 2.52, 1.18, "名目GDP\n実質GDP", size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.12), Inches(3.08), Inches(0.62), Inches(0.82))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(CYAN)
    arrow.line.fill.background()
    add_box(slide, 4.98, 2.15, 4.15, 2.70, "", fill="E8F1FF", line=BLUE)
    add_text(slide, 5.24, 2.40, 3.62, 0.36, "FORMULA", size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 5.22, 3.13, 3.65, 0.90, "GDPデフレーター\n＝ 名目GDP ÷ 実質GDP × 100", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(9.38), Inches(3.08), Inches(0.62), Inches(0.82))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(TEAL)
    arrow.line.fill.background()
    add_box(slide, 10.24, 2.15, 2.35, 2.70, "", fill="D9FBFB", line=TEAL)
    add_text(slide, 10.50, 2.40, 1.82, 0.36, "OUTPUT", size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 10.50, 3.26, 1.82, 0.60, "物価水準\n（指数）", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(
        slide,
        "記号を並べるだけでなく、何を入れると何が分かるかを図示する。",
        len(prs.slides),
    )


def pattern_kpi(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "G12",
        "KPI・表・指標比較",
        "重要数値はカードと小さなグラフを組み合わせ、比較の基準を明示する",
    )
    kpis = [
        ("実質GDP", "+1.2%", "前年比", BLUE),
        ("完全失業率", "2.6%", "低下＝改善傾向", TEAL),
        ("CPI", "+2.4%", "前年同月比", ORANGE),
    ]
    for i, (label, value, note, color) in enumerate(kpis):
        x = 0.66 + i * 3.02
        add_box(slide, x, 1.82, 2.72, 1.46, "", fill=WHITE, line=GRAY_20)
        add_text(slide, x + 0.18, 2.02, 2.34, 0.28, label, size=12, color=GRAY_70, bold=True)
        add_text(slide, x + 0.18, 2.38, 2.34, 0.46, value, size=25, color=color, bold=True)
        add_text(slide, x + 0.18, 2.88, 2.34, 0.22, note, size=9.5, color=GRAY_50)
    add_box(slide, 9.72, 1.82, 2.92, 3.88, "", fill=GRAY_10, line=GRAY_20)
    values = [42, 58, 49, 72]
    labels = ["Q1", "Q2", "Q3", "Q4"]
    for i, (value, label) in enumerate(zip(values, labels)):
        x = 10.08 + i * 0.58
        bar_h = value / 36
        add_box(slide, x, 5.14 - bar_h, 0.38, bar_h, "", fill=BLUE if i < 3 else TEAL, line=BLUE if i < 3 else TEAL, radius=False)
        add_text(slide, x - 0.04, 5.22, 0.46, 0.24, label, size=9, color=GRAY_70, align=PP_ALIGN.CENTER)
    table_rows = [
        ("指標", "見るポイント", "注意"),
        ("GDP", "実質成長率", "名目と区別"),
        ("失業率", "労働力人口も確認", "意欲喪失効果"),
        ("物価", "指数の対象範囲", "算式の違い"),
    ]
    for r, row in enumerate(table_rows):
        y = 3.62 + r * 0.52
        for c, value in enumerate(row):
            widths = [1.38, 3.05, 3.98]
            x = 0.66 + sum(widths[:c])
            add_box(
                slide,
                x,
                y,
                widths[c],
                0.46,
                value,
                fill=NAVY if r == 0 else (GRAY_10 if r % 2 else WHITE),
                line=WHITE if r == 0 else GRAY_20,
                size=10.5,
                color=WHITE if r == 0 else INK,
                bold=r == 0 or c == 0,
                align=PP_ALIGN.LEFT if c else PP_ALIGN.CENTER,
                radius=False,
            )
    add_takeaway(
        slide,
        "数値カードは3〜4個、表は4列×6行程度まで。強調色は1つに絞る。",
        len(prs.slides),
    )


def add_axes(slide, x: float, y: float, w: float, h: float, x_label: str, y_label: str) -> None:
    add_line(slide, x, y + h, x + w, y + h, color=NAVY, width=1.6)
    add_line(slide, x, y + h, x, y, color=NAVY, width=1.6)
    add_text(slide, x + w - 0.90, y + h + 0.12, 0.95, 0.24, x_label, size=10.5, color=NAVY, align=PP_ALIGN.RIGHT)
    add_text(slide, x - 0.16, y - 0.30, 1.05, 0.24, y_label, size=10.5, color=NAVY)


def econ_supply_demand(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "E01",
        "需要・供給曲線",
        "需要と供給の交点で価格と数量が同時に決まる",
        accent=TEAL,
    )
    x, y, w, h = 1.10, 1.90, 6.55, 3.95
    add_axes(slide, x, y, w, h, "数量 Q", "価格 P")
    add_line(slide, x + 0.35, y + 0.35, x + w - 0.35, y + h - 0.35, color=BLUE, width=3)
    add_line(slide, x + 0.35, y + h - 0.35, x + w - 0.35, y + 0.35, color=RED, width=3)
    add_text(slide, x + 5.80, y + 3.05, 0.50, 0.28, "S", size=14, color=BLUE, bold=True)
    add_text(slide, x + 5.80, y + 0.62, 0.50, 0.28, "D", size=14, color=RED, bold=True)
    eq_x, eq_y = x + w / 2, y + h / 2
    add_circle(slide, eq_x - 0.10, eq_y - 0.10, 0.20, "", fill=NAVY)
    add_line(slide, eq_x, eq_y, eq_x, y + h, color=GRAY_50, width=1)
    add_line(slide, x, eq_y, eq_x, eq_y, color=GRAY_50, width=1)
    add_box(slide, 8.30, 2.08, 4.00, 0.80, "需要 D：価格↑ → 需要量↓", fill="FFF1F1", line=RED, size=15, color=RED, bold=True)
    add_box(slide, 8.30, 3.18, 4.00, 0.80, "供給 S：価格↑ → 供給量↑", fill="E8F1FF", line=BLUE, size=15, color=BLUE, bold=True)
    add_box(slide, 8.30, 4.28, 4.00, 0.80, "交点 E：均衡価格・均衡数量", fill=GRAY_10, line=NAVY, size=15, color=NAVY, bold=True)
    add_takeaway(
        slide,
        "軸、曲線名、交点、補助線の4点を必ず表示する。",
        len(prs.slides),
        accent=TEAL,
    )


def econ_shift(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "E02",
        "曲線シフトと均衡変化",
        "需要増加は需要曲線を右へ動かし、均衡価格と均衡数量をともに押し上げる",
        accent=TEAL,
    )
    x, y, w, h = 0.95, 1.85, 7.10, 4.05
    add_axes(slide, x, y, w, h, "数量 Q", "価格 P")
    add_line(slide, x + 0.55, y + 0.30, x + w - 0.45, y + h - 0.40, color=BLUE, width=2.8)
    add_line(slide, x + 0.35, y + h - 0.55, x + w - 1.10, y + 0.40, color=GRAY_50, width=2.2)
    add_line(slide, x + 1.15, y + h - 0.55, x + w - 0.30, y + 0.40, color=RED, width=3)
    add_text(slide, x + 5.78, y + 3.15, 0.44, 0.24, "S", size=13, color=BLUE, bold=True)
    add_text(slide, x + 5.15, y + 0.62, 0.62, 0.24, "D₀", size=13, color=GRAY_50, bold=True)
    add_text(slide, x + 6.05, y + 0.62, 0.62, 0.24, "D₁", size=13, color=RED, bold=True)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.18), Inches(2.00), Inches(1.02), Inches(0.40))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(RED)
    arrow.line.fill.background()
    add_box(slide, 8.58, 2.02, 3.70, 0.70, "原因｜所得増加・嗜好変化など", fill=GRAY_10, line=GRAY_20, size=13, color=NAVY, bold=True)
    add_box(slide, 8.58, 3.10, 3.70, 0.70, "需要曲線｜D₀ → D₁（右方）", fill="FFF1F1", line=RED, size=13, color=RED, bold=True)
    add_box(slide, 8.58, 4.18, 3.70, 0.70, "結果｜価格↑・数量↑", fill="D9FBFB", line=TEAL, size=15, color=TEAL, bold=True)
    add_takeaway(
        slide,
        "曲線上の移動と曲線自体のシフトを、色と矢印で区別する。",
        len(prs.slides),
        accent=TEAL,
    )


def econ_keynesian_cross(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "E03",
        "45度線分析",
        "総需要が45度線と交わる点で均衡GDPが決まり、支出増加は乗数効果を生む",
        accent=TEAL,
    )
    x, y, w, h = 0.95, 1.90, 7.20, 3.95
    add_axes(slide, x, y, w, h, "GDP（所得）Y", "総需要 AE")
    add_line(slide, x, y + h, x + w - 0.40, y + 0.40, color=GRAY_50, width=2)
    add_text(slide, x + 5.82, y + 0.40, 0.90, 0.28, "45度線", size=11, color=GRAY_50, bold=True)
    add_line(slide, x + 0.35, y + 3.05, x + w - 0.45, y + 1.22, color=BLUE, width=3)
    add_line(slide, x + 0.35, y + 2.40, x + w - 0.45, y + 0.57, color=TEAL, width=3)
    add_text(slide, x + 5.94, y + 1.12, 0.68, 0.24, "AE₀", size=12, color=BLUE, bold=True)
    add_text(slide, x + 5.94, y + 0.47, 0.68, 0.24, "AE₁", size=12, color=TEAL, bold=True)
    add_box(slide, 8.60, 2.00, 3.60, 0.72, "政府支出 G が増加", fill="D9FBFB", line=TEAL, size=15, color=TEAL, bold=True)
    add_box(slide, 8.60, 3.08, 3.60, 0.72, "総需要線 AE が上方シフト", fill="E8F1FF", line=BLUE, size=14, color=BLUE, bold=True)
    add_box(slide, 8.60, 4.16, 3.60, 0.96, "均衡GDPの増加\nΔY ＝ 1/(1−c) × ΔG", fill=GRAY_10, line=NAVY, size=16, color=NAVY, bold=True)
    add_takeaway(
        slide,
        "45度線は支出＝所得を表す。シフト前後の均衡点を比較する。",
        len(prs.slides),
        accent=TEAL,
    )


def econ_is_lm(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "E04",
        "IS-LM分析",
        "財政政策はISを右へ、金融政策はLMを右へ動かし、GDPと利子率を変える",
        accent=TEAL,
    )
    x, y, w, h = 0.92, 1.90, 7.25, 3.95
    add_axes(slide, x, y, w, h, "GDP Y", "利子率 r")
    add_line(slide, x + 0.45, y + 0.40, x + w - 0.45, y + h - 0.40, color=RED, width=3)
    add_line(slide, x + 0.45, y + h - 0.40, x + w - 0.45, y + 0.40, color=BLUE, width=3)
    add_line(slide, x + 1.15, y + h - 0.40, x + w - 0.05, y + 0.40, color=TEAL, width=2.4)
    add_text(slide, x + 5.95, y + 3.04, 0.70, 0.24, "LM", size=12, color=RED, bold=True)
    add_text(slide, x + 5.92, y + 0.55, 0.70, 0.24, "IS₀", size=12, color=BLUE, bold=True)
    add_text(slide, x + 6.58, y + 0.55, 0.70, 0.24, "IS₁", size=12, color=TEAL, bold=True)
    add_box(slide, 8.55, 1.95, 3.75, 0.82, "財政政策：IS → 右", fill="D9FBFB", line=TEAL, size=15, color=TEAL, bold=True)
    add_box(slide, 8.55, 3.05, 3.75, 0.82, "結果：GDP ↑・利子率 ↑", fill=GRAY_10, line=GRAY_20, size=15, color=NAVY, bold=True)
    add_box(slide, 8.55, 4.15, 3.75, 1.02, "注意：利子率上昇が\n民間投資を押し出す", fill="FFF1F1", line=RED, size=14, color=RED, bold=True)
    add_takeaway(
        slide,
        "政策→曲線→交点→Y・rの順に追う。傾きと政策効果は別スライドにする。",
        len(prs.slides),
        accent=TEAL,
    )


def econ_surplus(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "E05",
        "余剰・死荷重",
        "課税で取引量が減ると、税収にならない余剰が死荷重として失われる",
        accent=TEAL,
    )
    x, y, w, h = 0.95, 1.86, 7.15, 4.00
    add_axes(slide, x, y, w, h, "数量 Q", "価格 P")
    add_line(slide, x + 0.35, y + 0.35, x + w - 0.35, y + h - 0.35, color=BLUE, width=2.8)
    add_line(slide, x + 0.35, y + h - 0.35, x + w - 0.35, y + 0.35, color=RED, width=2.8)
    # 余剰領域を半透明の代わりに淡色三角形で表現
    tri1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.12), Inches(2.15), Inches(2.95), Inches(1.60))
    tri1.rotation = 180
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = rgb("D0E2FF")
    tri1.line.color.rgb = rgb(BLUE)
    tri2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.12), Inches(3.75), Inches(2.95), Inches(1.60))
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = rgb("A7F0BA")
    tri2.line.color.rgb = rgb(GREEN)
    add_text(slide, 3.72, 2.74, 1.75, 0.30, "消費者余剰", size=12, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 3.72, 4.46, 1.75, 0.30, "生産者余剰", size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 8.52, 2.00, 3.78, 0.78, "税の導入 → 取引量減少", fill=GRAY_10, line=GRAY_20, size=14, color=NAVY, bold=True)
    add_box(slide, 8.52, 3.10, 3.78, 0.78, "余剰の一部 → 政府税収", fill="FFF8E1", line=YELLOW, size=14, color=NAVY, bold=True)
    add_box(slide, 8.52, 4.20, 3.78, 0.90, "残り → 死荷重\n社会全体から消失", fill="FFF1F1", line=RED, size=15, color=RED, bold=True)
    add_takeaway(
        slide,
        "消費者余剰・生産者余剰・税収・死荷重は色を固定して使う。",
        len(prs.slides),
        accent=TEAL,
    )


def econ_cycle_chart(prs: Presentation) -> None:
    slide = add_base(prs)
    add_header(
        slide,
        "E06",
        "景気循環・時系列",
        "トレンドと循環変動を分けて描くと、景気局面と転換点を判別できる",
        accent=TEAL,
    )
    x, y, w, h = 0.92, 1.82, 8.25, 4.08
    add_axes(slide, x, y, w, h, "時間 t", "実質GDP")
    add_line(slide, x + 0.20, y + h - 0.35, x + w - 0.30, y + 0.55, color=GRAY_50, width=1.8)
    points = [
        (x + 0.25, y + 3.45),
        (x + 1.30, y + 2.55),
        (x + 2.35, y + 1.55),
        (x + 3.30, y + 2.25),
        (x + 4.25, y + 2.95),
        (x + 5.25, y + 1.82),
        (x + 6.25, y + 0.88),
        (x + 7.20, y + 1.42),
        (x + 7.85, y + 1.92),
    ]
    for first, second in zip(points, points[1:]):
        add_line(slide, first[0], first[1], second[0], second[1], color=BLUE, width=3)
    add_text(slide, x + 2.00, y + 1.05, 0.72, 0.28, "山", size=13, color=RED, bold=True)
    add_text(slide, x + 4.00, y + 3.02, 0.72, 0.28, "谷", size=13, color=TEAL, bold=True)
    stages = [
        ("回復", "生産・雇用↑", BLUE),
        ("好況", "需要増・物価↑", GREEN),
        ("後退", "在庫増・投資↓", ORANGE),
        ("不況", "失業↑・政策対応", PURPLE),
    ]
    for i, (label, body, color) in enumerate(stages):
        y_box = 1.90 + i * 1.00
        add_box(slide, 9.58, y_box, 2.70, 0.76, "", fill=WHITE, line=color)
        add_box(slide, 9.70, y_box + 0.12, 0.75, 0.50, label, fill=color, line=color, size=12, color=WHITE, bold=True)
        add_text(slide, 10.62, y_box + 0.12, 1.42, 0.50, body, size=11.5, color=NAVY, bold=True)
    add_takeaway(
        slide,
        "実績線、長期トレンド、山・谷、4局面を同じグラフ上で区別する。",
        len(prs.slides),
        accent=TEAL,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "SMEC 標準図解スライドテンプレート"
    prs.core_properties.subject = "汎用図解12種＋経済学専用グラフ6種"
    prs.core_properties.author = "SMEC Project"
    prs.core_properties.keywords = "SMEC, 中小企業診断士, PowerPoint, 図解, テンプレート"

    add_title_slide(prs)
    add_design_rules(prs)
    add_selector(prs)

    generic_patterns = [
        pattern_definition,
        pattern_comparison,
        pattern_process,
        pattern_cause_effect,
        pattern_cycle,
        pattern_hierarchy,
        pattern_timeline,
        pattern_matrix,
        pattern_tree,
        pattern_network,
        pattern_formula,
        pattern_kpi,
    ]
    economics_patterns = [
        econ_supply_demand,
        econ_shift,
        econ_keynesian_cross,
        econ_is_lm,
        econ_surplus,
        econ_cycle_chart,
    ]
    for pattern in generic_patterns + economics_patterns:
        pattern(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Generated {len(prs.slides)} slides → {output}")


if __name__ == "__main__":
    build_deck(OUTPUT)
